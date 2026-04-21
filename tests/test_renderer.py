from __future__ import annotations

import base64
import io
import shutil
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

import pytest
from pptx import Presentation

from markdown_slides.assets import default_template_path
from markdown_slides.errors import TemplateError
from markdown_slides.parser import parse_deck
from markdown_slides.renderer import Downloader, render_pptx

PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9WHZ1x0AAAAASUVORK5CYII="
)
NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
TABLE_STYLE_MEDIUM_1_ACCENT_1 = "{B301B821-A1FF-4177-AEE7-76D212191A09}"


class FakeDownloader(Downloader):
    def fetch(self, url: str) -> bytes:
        return PNG_BYTES


def test_render_preserves_text_placeholders_and_notes(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
fonts:
  body: Aptos
  headings: Aptos Display
title_color: "#112233"
body_color: "#445566"
color_scheme:
  preset: Office
---

# Title slide title
---
layout: Title Slide
body_color: "#778899"
notes: |
  Note line.
---

Title slide subtitle

# Title and content title

Title and content text

# Section header title
---
layout: Section Header
---

Section header subtitle
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    prs = Presentation(str(output))
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text == "Title slide title"
    assert prs.slides[0].placeholders[1].text == "Title slide subtitle"
    assert prs.slides[1].shapes.title.text == "Title and content title"
    assert "Title and content text" in prs.slides[1].placeholders[1].text
    assert prs.slides[2].shapes.title.text == "Section header title"
    assert prs.slides[2].placeholders[1].text == "Section header subtitle"
    assert prs.slides[0].notes_slide.notes_text_frame.text == "Note line."
    with zipfile.ZipFile(output) as zf:
        slide1_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
        slide2_xml = zf.read("ppt/slides/slide2.xml").decode("utf-8")
    assert 'val="112233"' in slide1_xml
    assert 'val="778899"' in slide1_xml
    assert 'val="112233"' in slide2_xml
    assert 'val="445566"' in slide2_xml


def test_render_sets_theme_and_aspect_ratio(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
aspect_ratio: "4:3"
fonts:
  body: Aptos
  headings: Aptos Display
color_scheme:
  preset: Blue Warm
---

# Slide

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    prs = Presentation(str(output))
    assert prs.slide_width == 9144000
    assert prs.slide_height == 6858000
    with zipfile.ZipFile(output) as zf:
        theme = ET.fromstring(zf.read("ppt/theme/theme1.xml"))
    clr = theme.find(".//a:clrScheme", NS)
    assert clr is not None
    assert clr.attrib["name"] == "Blue Warm"
    accent1 = clr.find("a:accent1", NS)[0]
    assert accent1.attrib["val"] == "4A66AC"


def test_render_preserves_custom_template_theme_when_markdown_has_no_overrides(tmp_path: Path) -> None:
    template = tmp_path / "template.pptx"
    shutil.copyfile(default_template_path(), template)
    with zipfile.ZipFile(template, "r") as source, zipfile.ZipFile(
        tmp_path / "template-customized.pptx", "w", compression=zipfile.ZIP_DEFLATED
    ) as target:
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/theme/theme1.xml":
                root = ET.fromstring(data)
                clr = root.find(".//a:clrScheme", NS)
                assert clr is not None
                clr.set("name", "Custom Template Theme")
                accent1 = clr.find("a:accent1", NS)
                assert accent1 is not None
                accent1[0].set("val", "ABCDEF")
                font_scheme = root.find(".//a:fontScheme", NS)
                assert font_scheme is not None
                major = font_scheme.find("a:majorFont/a:latin", NS)
                minor = font_scheme.find("a:minorFont/a:latin", NS)
                assert major is not None
                assert minor is not None
                major.set("typeface", "Template Headings")
                minor.set("typeface", "Template Body")
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    customized_template = tmp_path / "template-customized.pptx"
    deck = parse_deck(
        "# Slide\n\nBody\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=customized_template, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        theme = ET.fromstring(zf.read("ppt/theme/theme1.xml"))
    clr = theme.find(".//a:clrScheme", NS)
    assert clr is not None
    assert clr.attrib["name"] == "Custom Template Theme"
    accent1 = clr.find("a:accent1", NS)
    assert accent1 is not None
    assert accent1[0].attrib["val"] == "ABCDEF"
    font_scheme = theme.find(".//a:fontScheme", NS)
    assert font_scheme is not None
    major = font_scheme.find("a:majorFont/a:latin", NS)
    minor = font_scheme.find("a:minorFont/a:latin", NS)
    assert major is not None
    assert minor is not None
    assert major.attrib["typeface"] == "Template Headings"
    assert minor.attrib["typeface"] == "Template Body"


def test_render_background_and_body_image(tmp_path: Path) -> None:
    image_path = tmp_path / "photo.png"
    image_path.write_bytes(PNG_BYTES)
    deck = parse_deck(
        """---
background: "linear-gradient(90deg, #0E2841 0%, #156082 100%)"
color_scheme:
  preset: Office
---

# Photo

![alt](./photo.png)
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    prs = Presentation(str(output))
    names = [shape.name for shape in prs.slides[0].shapes]
    assert prs.slide_masters[0].background.fill.type == 3
    assert round(prs.slide_masters[0].background.fill.gradient_angle) == 90
    assert "MarkdownSlidesImage" in names


def test_render_remote_image_with_fake_downloader(tmp_path: Path) -> None:
    deck = parse_deck(
        "# Slide\n\n![alt](https://example.com/photo.png)\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(
        deck,
        output_path=output,
        template_path=None,
        force=False,
        base_dir=tmp_path,
        downloader=FakeDownloader(),
    )

    prs = Presentation(str(output))
    assert "MarkdownSlidesImage" in [shape.name for shape in prs.slides[0].shapes]


def test_render_lists_and_headings_have_expected_bullet_xml(tmp_path: Path) -> None:
    deck = parse_deck(
        "# Slide\n\n## Heading\nParagraph.\n\n- Bullet\n1. Numbered\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:buNone/>" in slide_xml
    assert "arabicPeriod" in slide_xml
    assert "• Bullet" not in slide_xml


def test_render_document_background_image_targets_slide_master(tmp_path: Path) -> None:
    image_path = tmp_path / "bg.png"
    image_path.write_bytes(PNG_BYTES)
    deck = parse_deck(
        """---
background: "url('./bg.png')"
---

# Slide

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        master_xml = zf.read("ppt/slideMasters/slideMaster1.xml").decode("utf-8")
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:blipFill>" in master_xml
    assert 'r:embed="' in master_xml
    assert "MarkdownSlidesBackgroundImage" not in slide_xml


def test_render_theme_font_refs_blockquote_indent_and_table_style(tmp_path: Path) -> None:
    deck = parse_deck(
        """# Formatting

## Heading
Paragraph

> Quote text

# Table

| A | B |
| --- | --- |
| 1 | 2 |
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
        table_slide_xml = zf.read("ppt/slides/slide2.xml").decode("utf-8")
    assert '+mj-lt' in slide_xml
    assert '+mn-lt' in slide_xml
    assert 'marL="0"' in slide_xml
    assert 'indent="0"' in slide_xml
    assert 'schemeClr val="accent1"' in slide_xml
    assert TABLE_STYLE_MEDIUM_1_ACCENT_1 in table_slide_xml


def test_render_radial_gradient_background_writes_path_gradient_xml(tmp_path: Path) -> None:
    deck = parse_deck(
        """# Radial
---
background: "radial-gradient(circle, #0E2841 0%, #156082 55%, #EAF3FF 100%)"
---

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert '<a:path path="circle">' in slide_xml
    assert '<a:fillToRect l="50000" t="50000" r="50000" b="50000"/>' in slide_xml


def test_render_theme_color_refs_for_text_and_backgrounds(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
title_color: "var(--light-1)"
body_color: "var(--dark-1)"
background: "linear-gradient(90deg, var(--accent-1) 0%, var(--accent-2) 100%)"
---

# Slide

Paragraph

> Quote text
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        master_xml = zf.read("ppt/slideMasters/slideMaster1.xml").decode("utf-8")
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'schemeClr val="accent1"' in master_xml
    assert 'schemeClr val="accent2"' in master_xml
    assert 'schemeClr val="lt1"' in slide_xml
    assert 'schemeClr val="dk1"' in slide_xml
    assert 'schemeClr val="accent1"' in slide_xml


def test_render_section_header_subtitle_uses_theme_body_text_color_by_default(tmp_path: Path) -> None:
    deck = parse_deck(
        """# Section heading
---
layout: Section Header
---

Subtitle text
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'schemeClr val="dk1"' in slide_xml
    assert "<a:defRPr" in slide_xml
    assert "<a:endParaRPr" in slide_xml


def test_force_flag_controls_overwrite(tmp_path: Path) -> None:
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))
    output = tmp_path / "deck.pptx"
    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with pytest.raises(Exception):
        render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    render_pptx(deck, output_path=output, template_path=None, force=True, base_dir=tmp_path)
    assert output.exists()


def test_missing_placeholder_template_fails(tmp_path: Path) -> None:
    template = tmp_path / "broken-template.pptx"
    with zipfile.ZipFile(default_template_path(), "r") as source, zipfile.ZipFile(template, "w", compression=zipfile.ZIP_DEFLATED) as target:
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/slideLayouts/slideLayout2.xml":
                xml = ET.fromstring(data)
                sp_tree = xml.find(".//p:spTree", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
                for sp in list(sp_tree):
                    ph = sp.find(".//p:ph", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
                    if ph is not None and ph.attrib.get("idx") == "1":
                        sp_tree.remove(sp)
                        break
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))
    output = tmp_path / "deck.pptx"

    with pytest.raises(TemplateError):
        render_pptx(deck, output_path=output, template_path=template, force=False, base_dir=tmp_path)
