from __future__ import annotations

import io
import re
import shutil
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
import zipfile

import httpx
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.base import BaseShape
from pptx.util import Emu, Pt

from markdown_slides.assets import default_template_path
from markdown_slides.errors import AssetError, RenderError, TemplateError
from markdown_slides.models import Background, BodyContent, Deck, InlineText, Slide, TableBlock

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
TITLE_PLACEHOLDERS = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
SUBTITLE_PLACEHOLDERS = {PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY}
BODY_PLACEHOLDERS = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
TABLE_STYLE_MEDIUM_1_ACCENT_1 = "{B301B821-A1FF-4177-AEE7-76D212191A09}"
HEADING_SCALE = {2: 1.33, 3: 1.2, 4: 1.1, 5: 1.05, 6: 1.0}
DEFAULT_BODY_LINE_SPACING = 1.0
DEFAULT_BODY_SPACE_BEFORE_PT = 12
DEFAULT_BODY_SPACE_AFTER_PT = 6
THEME_COLOR_VAR_RE = re.compile(r"^var\(\s*--(?P<name>[a-z0-9-]+)\s*\)$", re.IGNORECASE)
THEME_COLOR_SCHEME_MAP = {
    "dark-1": "dk1",
    "light-1": "lt1",
    "dark-2": "dk2",
    "light-2": "lt2",
    "accent-1": "accent1",
    "accent-2": "accent2",
    "accent-3": "accent3",
    "accent-4": "accent4",
    "accent-5": "accent5",
    "accent-6": "accent6",
    "hyperlink": "hlink",
    "followed-hyperlink": "folHlink",
}
def render_pptx(
    deck: Deck,
    *,
    output_path: Path,
    template_path: Path | None,
    force: bool,
    base_dir: Path,
    downloader: Downloader | None = None,
) -> Path:
    if output_path.exists() and not force:
        raise RenderError("output_exists", f"Output file already exists: {output_path}")
    template = template_path or default_template_path()
    preserve_template_paragraph_formatting = template_path is not None
    presentation = Presentation(str(template))
    _clear_existing_slides(presentation)
    _apply_aspect_ratio(presentation, deck.aspect_ratio)
    template_defaults = _read_template_defaults(presentation)
    if deck.background is not None:
        _apply_master_background(presentation.slide_masters[0], deck.background, base_dir=base_dir, downloader=downloader or Downloader())
    downloader = downloader or Downloader()
    layout_map = {layout.name.casefold(): layout for layout in presentation.slide_layouts}

    for slide_spec in deck.slides:
        layout = layout_map.get(slide_spec.layout.casefold())
        if layout is None:
            raise TemplateError("missing_layout", f"Template does not contain layout '{slide_spec.layout}'.")
        slide = presentation.slides.add_slide(layout)
        _apply_hide_background_graphics(slide, slide_spec)
        if slide_spec.background is not None:
            _apply_background(
                slide,
                slide_spec.background,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
                base_dir=base_dir,
                downloader=downloader,
            )
        _render_title(slide, slide_spec, deck)
        _render_body(
            slide,
            slide_spec,
            deck,
            template_defaults=template_defaults,
            preserve_template_paragraph_formatting=preserve_template_paragraph_formatting,
            base_dir=base_dir,
            downloader=downloader,
        )
        _render_notes(slide, slide_spec)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
        temp_path = Path(temp_file.name)
    try:
        presentation.save(str(temp_path))
        _rewrite_theme(temp_path, deck)
        shutil.move(str(temp_path), str(output_path))
    finally:
        if temp_path.exists():
            temp_path.unlink(missing_ok=True)
    return output_path


class Downloader:
    def fetch(self, url: str) -> bytes:
        response = httpx.get(url, follow_redirects=True, timeout=30.0)
        response.raise_for_status()
        return response.content


def list_layouts(template_path: Path | None = None) -> list[str]:
    presentation = Presentation(str(template_path or default_template_path()))
    return [layout.name for layout in presentation.slide_layouts]


def _clear_existing_slides(presentation: Presentation) -> None:
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        slide_id_list.remove(slide_id)


def _apply_aspect_ratio(presentation: Presentation, aspect_ratio: str) -> None:
    if aspect_ratio == "16:9":
        presentation.slide_width = Emu(12192000)
        presentation.slide_height = Emu(6858000)
    else:
        presentation.slide_width = Emu(9144000)
        presentation.slide_height = Emu(6858000)


def _render_title(slide, slide_spec: Slide, deck: Deck) -> None:
    if slide_spec.layout == "Blank":
        return
    title_shape = _require_placeholder(slide, TITLE_PLACEHOLDERS, "title")
    text_frame = title_shape.text_frame
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = slide_spec.title
    title_color = _resolve_text_color(deck, slide_spec, "title")
    if title_color is not None:
        _set_paragraph_default_color(paragraph, title_color)
        _set_run_color(run, title_color)


def _render_body(
    slide,
    slide_spec: Slide,
    deck: Deck,
    *,
    template_defaults: dict[str, float],
    preserve_template_paragraph_formatting: bool,
    base_dir: Path,
    downloader: Downloader,
) -> None:
    body = slide_spec.body
    body_color = _resolve_text_color(deck, slide_spec, "body")
    if body_color is None and slide_spec.layout in {"Title Slide", "Section Header"}:
        body_color = "var(--dark-1)"
    if slide_spec.layout in {"Blank", "Title Only"} or body.is_empty:
        return
    if slide_spec.layout in {"Title Slide", "Section Header"}:
        placeholder = _require_placeholder(slide, SUBTITLE_PLACEHOLDERS, "subtitle")
        _render_text_flow(
            placeholder,
            body,
            deck,
            template_defaults=template_defaults,
            text_color=body_color,
            preserve_template_paragraph_formatting=preserve_template_paragraph_formatting,
        )
        return
    if slide_spec.layout != "Title and Content":
        raise TemplateError("unsupported_layout", f"Layout '{slide_spec.layout}' is not renderable.")
    placeholder = _require_placeholder(slide, BODY_PLACEHOLDERS, "body")
    if body.paragraphs:
        _render_text_flow(
            placeholder,
            body,
            deck,
            template_defaults=template_defaults,
            text_color=body_color,
            preserve_template_paragraph_formatting=preserve_template_paragraph_formatting,
        )
        return
    if body.images:
        _render_image(slide, placeholder, body.images[0].src, base_dir=base_dir, downloader=downloader, contain=True, name="MarkdownSlidesImage")
        return
    if body.tables:
        _render_table(slide, placeholder, body.tables[0], deck)
        return


def _render_text_flow(
    placeholder,
    body: BodyContent,
    deck: Deck,
    *,
    template_defaults: dict[str, float],
    text_color: str | None,
    preserve_template_paragraph_formatting: bool,
) -> None:
    text_frame = placeholder.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraphs = body.paragraphs
    for index, paragraph_model in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.clear()
        _configure_paragraph_bullets(paragraph, paragraph_model)
        _apply_paragraph_spacing(paragraph, paragraph_model, preserve_template_paragraph_formatting=preserve_template_paragraph_formatting)
        if text_color is not None and paragraph_model.kind != "blockquote":
            _set_paragraph_default_color(paragraph, text_color)
        if paragraph_model.kind == "blockquote":
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(6)
        fragments = paragraph_model.fragments or [InlineText(kind="text", text="")]
        for fragment in fragments:
            _add_fragment_runs(paragraph, fragment, deck, paragraph_model, template_defaults=template_defaults, text_color=text_color)
        if not paragraph.runs:
            run = paragraph.add_run()
            run.text = ""


def _add_fragment_runs(paragraph, fragment: InlineText, deck: Deck, paragraph_model, *, template_defaults: dict[str, float], text_color: str | None) -> None:
    if fragment.kind == "link":
        for child in fragment.children:
            run = paragraph.add_run()
            run.text = _flatten_inline([child])
            run.hyperlink.address = fragment.href
            _apply_run_font(run, deck, paragraph_model, fragment.kind, template_defaults=template_defaults, text_color=text_color)
        return
    if fragment.children:
        for child in fragment.children:
            run = paragraph.add_run()
            run.text = _flatten_inline([child])
            _apply_run_font(run, deck, paragraph_model, fragment.kind, template_defaults=template_defaults, text_color=text_color)
        return
    run = paragraph.add_run()
    run.text = fragment.text or ""
    _apply_run_font(run, deck, paragraph_model, fragment.kind, template_defaults=template_defaults, text_color=text_color)


def _apply_run_font(run, deck: Deck, paragraph_model, inline_kind: str, *, template_defaults: dict[str, float], text_color: str | None) -> None:
    if paragraph_model.kind == "heading":
        _set_theme_font(run, major=True)
        base_size = template_defaults["body_font_pt"]
        scale = HEADING_SCALE.get(paragraph_model.heading_level or 6, 1.0)
        run.font.size = Pt(round(base_size * scale))
    elif paragraph_model.kind == "code" or inline_kind == "code":
        run.font.name = "Consolas"
    else:
        _set_theme_font(run, major=False)
    if paragraph_model.kind == "heading":
        run.font.bold = True
    elif paragraph_model.kind == "blockquote":
        run.font.italic = True
        _set_run_color(run, "var(--accent-1)")
    else:
        run.font.bold = inline_kind == "strong"
        run.font.italic = inline_kind == "emphasis"
    if inline_kind == "strong":
        run.font.bold = True
    if inline_kind == "emphasis":
        run.font.italic = True
    if text_color is not None and paragraph_model.kind != "blockquote":
        _set_run_color(run, text_color)


def _resolve_text_color(deck: Deck, slide_spec: Slide, kind: str) -> str | None:
    if slide_spec.text_colors is not None:
        value = getattr(slide_spec.text_colors, kind)
        if value is not None:
            return value
    if deck.text_colors is not None:
        return getattr(deck.text_colors, kind)
    return None


def _set_run_color(run, color_value: str) -> None:
    scheme_name = _theme_scheme_name(color_value)
    if scheme_name is None:
        run.font.color.rgb = RGBColor.from_string(color_value[1:])
        return
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag == qn("a:solidFill"):
            r_pr.remove(child)
    solid_fill = OxmlElement("a:solidFill")
    _append_color_choice(solid_fill, color_value)
    r_pr.append(solid_fill)


def _set_paragraph_default_color(paragraph, color_value: str) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    def_rpr = p_pr.get_or_add_defRPr()
    _set_text_character_color(def_rpr, color_value)
    end_rpr = paragraph._p.get_or_add_endParaRPr()
    _set_text_character_color(end_rpr, color_value)


def _set_text_character_color(r_pr, color_value: str) -> None:
    for child in list(r_pr):
        if child.tag == qn("a:solidFill"):
            r_pr.remove(child)
    solid_fill = OxmlElement("a:solidFill")
    _append_color_choice(solid_fill, color_value)
    r_pr.append(solid_fill)


def _apply_paragraph_spacing(paragraph, paragraph_model, *, preserve_template_paragraph_formatting: bool) -> None:
    if preserve_template_paragraph_formatting:
        return
    if paragraph_model.kind != "paragraph":
        return
    paragraph.line_spacing = DEFAULT_BODY_LINE_SPACING
    paragraph.space_before = Pt(DEFAULT_BODY_SPACE_BEFORE_PT)
    paragraph.space_after = Pt(DEFAULT_BODY_SPACE_AFTER_PT)


def _render_table(slide, placeholder, table: TableBlock, deck: Deck) -> None:
    rows = 1 + len(table.rows)
    cols = len(table.headers)
    if cols == 0:
        raise RenderError("invalid_table", "Tables must contain at least one header cell.")
    shape = slide.shapes.add_table(rows, cols, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
    shape.name = "MarkdownSlidesTable"
    table_shape = shape.table
    table_style_id = table_shape._tbl.tblPr.find(qn("a:tableStyleId"))
    if table_style_id is not None:
        table_style_id.text = TABLE_STYLE_MEDIUM_1_ACCENT_1
    for column_index, cell in enumerate(table.headers):
        table_shape.cell(0, column_index).text = _flatten_inline(cell)
    for row_index, row in enumerate(table.rows, start=1):
        for column_index, cell in enumerate(row):
            table_shape.cell(row_index, column_index).text = _flatten_inline(cell)
    for row in table_shape.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_theme_font(run, major=False)


def _render_image(slide, placeholder, src: str, *, base_dir: Path, downloader: Downloader, contain: bool, name: str) -> None:
    image_source = _resolve_image_source(src, base_dir=base_dir, downloader=downloader)
    left, top, width, height = _fit_image(image_source, placeholder.left, placeholder.top, placeholder.width, placeholder.height, contain=contain)
    picture = slide.shapes.add_picture(image_source, left, top, width, height)
    picture.name = name


def _resolve_image_source(src: str, *, base_dir: Path, downloader: Downloader):
    if src.startswith(("http://", "https://")):
        try:
            return io.BytesIO(downloader.fetch(src))
        except Exception as exc:  # noqa: BLE001
            raise AssetError("image_download_failed", f"Failed to download image '{src}': {exc}") from exc
    path = (base_dir / src).resolve()
    if not path.exists():
        raise AssetError("missing_asset", f"Image asset does not exist: {path}")
    return str(path)


def _fit_image(image_source, left: int, top: int, width: int, height: int, *, contain: bool) -> tuple[int, int, int, int]:
    if isinstance(image_source, io.BytesIO):
        image_source.seek(0)
        with PILImage.open(image_source) as image:
            image_width, image_height = image.size
        image_source.seek(0)
    else:
        with PILImage.open(image_source) as image:
            image_width, image_height = image.size
    scale = min(width / image_width, height / image_height) if contain else max(width / image_width, height / image_height)
    rendered_width = int(image_width * scale)
    rendered_height = int(image_height * scale)
    rendered_left = int(left + (width - rendered_width) / 2)
    rendered_top = int(top + (height - rendered_height) / 2)
    return rendered_left, rendered_top, rendered_width, rendered_height


def _render_notes(slide, slide_spec: Slide) -> None:
    if not slide_spec.notes:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = slide_spec.notes


def _apply_hide_background_graphics(slide, slide_spec: Slide) -> None:
    if slide_spec.hide_background_graphics:
        slide._element.set("showMasterSp", "0")


def _apply_background(slide, background: Background, *, slide_width: int, slide_height: int, base_dir: Path, downloader: Downloader) -> None:
    if background.kind == "none":
        slide.background.fill.background()
        return
    if background.kind == "color":
        if _theme_scheme_name(background.value or "") is not None:
            _apply_background_fill_xml(slide._element, background)
            return
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(background.value[1:])
        return
    if background.kind == "gradient":
        if background.gradient_kind == "radial" or _gradient_uses_theme_colors(background):
            _apply_background_fill_xml(slide._element, background)
            return
        fill = slide.background.fill
        fill.gradient()
        if background.angle is not None:
            fill.gradient_angle = background.angle
        stops = fill.gradient_stops
        for index, stop in enumerate(background.stops[: len(stops)]):
            stops[index].position = stop.position
            stops[index].color.rgb = RGBColor.from_string(stop.color[1:])
        return
    if background.kind == "image":
        image_source = _resolve_image_source(background.url or "", base_dir=base_dir, downloader=downloader)
        left, top, width, height = _fit_image(
            image_source,
            0,
            0,
            slide_width,
            slide_height,
            contain=False,
        )
        picture = slide.shapes.add_picture(image_source, left, top, width, height)
        picture.name = "MarkdownSlidesBackgroundImage"
        _send_to_back(slide, picture)
        return
    raise RenderError("invalid_background", f"Unsupported background kind '{background.kind}'.")


def _apply_master_background(master, background: Background, *, base_dir: Path, downloader: Downloader) -> None:
    if background.kind == "none":
        master.background.fill.background()
        return
    if background.kind == "color":
        if _theme_scheme_name(background.value or "") is not None:
            _apply_background_fill_xml(master._element, background)
            return
        fill = master.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(background.value[1:])
        return
    if background.kind == "gradient":
        if background.gradient_kind == "radial" or _gradient_uses_theme_colors(background):
            _apply_background_fill_xml(master._element, background)
            return
        fill = master.background.fill
        fill.gradient()
        if background.angle is not None:
            fill.gradient_angle = background.angle
        stops = fill.gradient_stops
        for index, stop in enumerate(background.stops[: len(stops)]):
            stops[index].position = stop.position
            stops[index].color.rgb = RGBColor.from_string(stop.color[1:])
        return
    if background.kind != "image":
        raise RenderError("invalid_background", f"Unsupported master background kind '{background.kind}'.")
    image_source = _resolve_image_source(background.url or "", base_dir=base_dir, downloader=downloader)
    _image_part, r_id = master.part.get_or_add_image_part(image_source)
    c_sld = master._element.find(qn("p:cSld"))
    if c_sld is None:
        raise RenderError("missing_master", "Slide master content is missing.")
    existing_bg = c_sld.find(qn("p:bg"))
    if existing_bg is not None:
        c_sld.remove(existing_bg)
    bg = OxmlElement("p:bg")
    bg_pr = OxmlElement("p:bgPr")
    blip_fill = OxmlElement("a:blipFill")
    blip = OxmlElement("a:blip")
    blip.set(qn("r:embed"), r_id)
    stretch = OxmlElement("a:stretch")
    fill_rect = OxmlElement("a:fillRect")
    stretch.append(fill_rect)
    blip_fill.append(blip)
    blip_fill.append(stretch)
    bg_pr.append(blip_fill)
    bg_pr.append(OxmlElement("a:effectLst"))
    bg.append(bg_pr)
    c_sld.insert(0, bg)


def _apply_background_fill_xml(container, background: Background) -> None:
    if background.kind == "color":
        solid_fill = OxmlElement("a:solidFill")
        _append_color_choice(solid_fill, background.value or "")
        _set_background_fill_xml(container, solid_fill)
        return
    if background.kind != "gradient":
        raise RenderError("invalid_background", f"Unsupported XML background kind '{background.kind}'.")
    grad_fill = OxmlElement("a:gradFill")
    grad_fill.set("rotWithShape", "1")
    gs_list = OxmlElement("a:gsLst")
    for stop in background.stops:
        gs = OxmlElement("a:gs")
        gs.set("pos", str(int(stop.position * 100000)))
        _append_color_choice(gs, stop.color)
        gs_list.append(gs)
    grad_fill.append(gs_list)
    if background.gradient_kind == "radial":
        path = OxmlElement("a:path")
        path.set("path", "circle")
        fill_to_rect = OxmlElement("a:fillToRect")
        fill_to_rect.set("l", "50000")
        fill_to_rect.set("t", "50000")
        fill_to_rect.set("r", "50000")
        fill_to_rect.set("b", "50000")
        path.append(fill_to_rect)
        grad_fill.append(path)
    else:
        lin = OxmlElement("a:lin")
        lin.set("ang", str(_gradient_angle_to_ooxml(background.angle or 180.0)))
        lin.set("scaled", "0")
        grad_fill.append(lin)
    _set_background_fill_xml(container, grad_fill)


def _set_background_fill_xml(container, fill_element) -> None:
    c_sld = container.find(qn("p:cSld"))
    if c_sld is None:
        raise RenderError("missing_background", "Slide background container is missing.")
    existing_bg = c_sld.find(qn("p:bg"))
    if existing_bg is not None:
        c_sld.remove(existing_bg)
    bg = OxmlElement("p:bg")
    bg_pr = OxmlElement("p:bgPr")
    bg_pr.append(fill_element)
    bg_pr.append(OxmlElement("a:effectLst"))
    bg.append(bg_pr)
    c_sld.insert(0, bg)


def _append_color_choice(parent, color_value: str):
    scheme_name = _theme_scheme_name(color_value)
    if scheme_name is not None:
        color = OxmlElement("a:schemeClr")
        color.set("val", scheme_name)
    else:
        color = OxmlElement("a:srgbClr")
        color.set("val", color_value[1:])
    parent.append(color)
    return color


def _theme_scheme_name(color_value: str) -> str | None:
    if color_value.startswith("scheme:"):
        return color_value.split(":", 1)[1]
    match = THEME_COLOR_VAR_RE.match(color_value.strip())
    if match is None:
        return None
    return THEME_COLOR_SCHEME_MAP.get(match.group("name").lower())


def _gradient_uses_theme_colors(background: Background) -> bool:
    return any(_theme_scheme_name(stop.color) is not None for stop in background.stops)


def _gradient_angle_to_ooxml(angle: float) -> int:
    normalized = angle % 360.0
    return int(((360.0 - normalized) % 360.0) * 60000)


def _send_to_back(slide, shape: BaseShape) -> None:
    tree = slide.shapes._spTree
    element = shape._element
    tree.remove(element)
    tree.insert(2, element)


def _require_placeholder(slide, allowed_types: set[PP_PLACEHOLDER], kind: str):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type in allowed_types:
            return placeholder
    raise TemplateError("missing_placeholder", f"Selected layout does not contain the required {kind} placeholder.")


def _flatten_inline(fragments: list[InlineText]) -> str:
    parts: list[str] = []
    for fragment in fragments:
        if fragment.children:
            parts.append(_flatten_inline(fragment.children))
        else:
            parts.append(fragment.text or "")
    return "".join(parts)


def _configure_paragraph_bullets(paragraph, paragraph_model) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    _remove_bullet_elements(p_pr)
    if paragraph_model.kind == "list_item":
        paragraph.level = paragraph_model.level
        _clear_indent(p_pr)
        if paragraph_model.ordered_index is not None:
            auto_num = OxmlElement("a:buAutoNum")
            auto_num.set("type", "arabicPeriod")
            auto_num.set("startAt", str(paragraph_model.ordered_index))
            p_pr.insert(0, auto_num)
        return
    paragraph.level = 0
    _set_indent_attrs(p_pr, mar_l=0, indent=0)
    bu_none = OxmlElement("a:buNone")
    p_pr.insert(0, bu_none)


def _remove_bullet_elements(p_pr) -> None:
    bullet_tags = {
        qn("a:buNone"),
        qn("a:buAutoNum"),
        qn("a:buChar"),
        qn("a:buBlip"),
        qn("a:buClr"),
        qn("a:buClrTx"),
        qn("a:buFont"),
        qn("a:buFontTx"),
        qn("a:buSzPct"),
        qn("a:buSzPts"),
        qn("a:buSzTx"),
    }
    for child in list(p_pr):
        if child.tag in bullet_tags:
            p_pr.remove(child)


def _set_indent_attrs(p_pr, *, mar_l: int | None = None, indent: int | None = None) -> None:
    if mar_l is not None:
        p_pr.set("marL", str(mar_l))
    if indent is not None:
        p_pr.set("indent", str(indent))


def _clear_indent(p_pr) -> None:
    if "marL" in p_pr.attrib:
        del p_pr.attrib["marL"]
    if "indent" in p_pr.attrib:
        del p_pr.attrib["indent"]


def _set_paragraph_indent(paragraph, *, left: float, hanging: float) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    _set_indent_attrs(
        p_pr,
        mar_l=int(left * 914400),
        indent=int(-hanging * 914400),
    )


def _set_theme_font(run, *, major: bool) -> None:
    if run.font.name:
        run.font.name = None
    r_pr = run._r.get_or_add_rPr()
    for child_tag in (qn("a:latin"), qn("a:ea"), qn("a:cs")):
        for child in list(r_pr):
            if child.tag == child_tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    ea = OxmlElement("a:ea")
    cs = OxmlElement("a:cs")
    if major:
        latin.set("typeface", "+mj-lt")
        ea.set("typeface", "+mj-ea")
        cs.set("typeface", "+mj-cs")
    else:
        latin.set("typeface", "+mn-lt")
        ea.set("typeface", "+mn-ea")
        cs.set("typeface", "+mn-cs")
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def _read_template_defaults(presentation: Presentation) -> dict[str, float]:
    master = presentation.slide_masters[0]._element
    body_def = master.find(".//p:txStyles/p:bodyStyle/a:lvl1pPr/a:defRPr", NS)
    title_def = master.find(".//p:txStyles/p:titleStyle/a:lvl1pPr/a:defRPr", NS)
    body_pt = _sz_to_pt(body_def.get("sz") if body_def is not None else None, default=28.0)
    title_pt = _sz_to_pt(title_def.get("sz") if title_def is not None else None, default=44.0)
    return {"body_font_pt": body_pt, "title_font_pt": title_pt}


def _sz_to_pt(value: str | None, *, default: float) -> float:
    if value is None:
        return default
    try:
        return int(value) / 100.0
    except ValueError:
        return default


def _rewrite_theme(path: Path, deck: Deck) -> None:
    temp_path = path.with_suffix(".rewritten.pptx")
    color_map = {
        "dk1": "dark_1",
        "lt1": "light_1",
        "dk2": "dark_2",
        "lt2": "light_2",
        "accent1": "accent_1",
        "accent2": "accent_2",
        "accent3": "accent_3",
        "accent4": "accent_4",
        "accent5": "accent_5",
        "accent6": "accent_6",
        "hlink": "hyperlink",
        "folHlink": "followed_hyperlink",
    }
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(temp_path, "w", compression=zipfile.ZIP_DEFLATED) as target:
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/theme/theme1.xml":
                root = ET.fromstring(data)
                if deck.color_scheme is not None:
                    clr_scheme = root.find(".//a:clrScheme", NS)
                    if clr_scheme is None:
                        raise RenderError("missing_theme", "The template does not contain a theme color scheme.")
                    clr_scheme.set("name", deck.color_scheme.name)
                    for xml_key, model_key in color_map.items():
                        parent = clr_scheme.find(f"a:{xml_key}", NS)
                        if parent is None or len(parent) == 0:
                            continue
                        child = parent[0]
                        if child.tag.endswith("sysClr"):
                            child.set("lastClr", deck.color_scheme.colors[model_key][1:])
                        else:
                            child.set("val", deck.color_scheme.colors[model_key][1:])
                font_scheme = root.find(".//a:fontScheme", NS)
                if font_scheme is not None and deck.fonts_override:
                    major = font_scheme.find("a:majorFont/a:latin", NS)
                    minor = font_scheme.find("a:minorFont/a:latin", NS)
                    if major is not None:
                        major.set("typeface", deck.fonts.headings)
                    if minor is not None:
                        minor.set("typeface", deck.fonts.body)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    temp_path.replace(path)
